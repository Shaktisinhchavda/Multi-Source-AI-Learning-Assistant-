"""
PPTX processor — extracts text from PowerPoint slides.
"""

from io import BytesIO
from pptx import Presentation
from .chunker import chunk_pages


def extract_pptx_slides(file_bytes: bytes) -> list[dict]:
    """
    Extract text from each slide of a PPTX file.
    Returns list of {text, ref} dicts.
    """
    prs = Presentation(BytesIO(file_bytes))
    slides = []

    for i, slide in enumerate(prs.slides):
        texts = []

        # Extract text from all shapes
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    para_text = paragraph.text.strip()
                    if para_text:
                        texts.append(para_text)

            # Also extract from tables
            if shape.has_table:
                for row in shape.table.rows:
                    row_texts = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_texts.append(cell_text)
                    if row_texts:
                        texts.append(" | ".join(row_texts))

        slide_text = "\n".join(texts)
        if slide_text.strip():
            slides.append({
                "text": slide_text,
                "ref": f"slide {i + 1}",
            })

    return slides


def process_pptx(file_bytes: bytes, filename: str) -> dict:
    """
    Full PPTX processing pipeline:
    1. Extract text from each slide
    2. Chunk with slide references
    3. Generate a brief summary

    Returns: {chunks: [...], summary: str, slide_count: int}
    """
    slides = extract_pptx_slides(file_bytes)

    if not slides:
        raise ValueError(
            f"Could not extract any text from '{filename}'. "
            "The presentation may contain only images."
        )

    chunks = chunk_pages(
        pages=slides,
        source_type="pptx",
        source_name=filename,
    )

    # Create summary from first 2 slides
    all_text = " ".join(s["text"] for s in slides[:2])
    summary = all_text[:500].strip()
    if len(all_text) > 500:
        summary += "..."

    return {
        "chunks": chunks,
        "summary": summary,
        "slide_count": len(slides),
    }
